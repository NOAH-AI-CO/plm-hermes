import time

def get_main_py_init_content() -> str:
    content = """
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptxsdk import *

# 如果你需要其他模块，请在这里进行 import 你只能 import python 标准库 和 pptx 和 pptxsdk，禁止 import 其他第三方模块


# !!! *** 此区域禁止修改 开始 *** !!!
OUTPUT_PATH = "myppt.pptx"
PPT_W = Inches(13.333)
PPT_H = Inches(7.5)  # W:H = 16:9
PRS = Presentation()
PRS.slide_width = PPT_W
PRS.slide_height = PPT_H
EMPTY_LAYOUT = PRS.slide_layouts[6]


def new_slide():
    return PRS.slides.add_slide(EMPTY_LAYOUT)
# !!! *** 此区域禁止修改 结束 *** !!!

TOTAL_PAGE_COUNT = 0 # 此 PPT 的总页数，你要如实填写，变量名必须为 TOTAL_PAGE_COUNT，这个变量会被其他人查看，方便核对页数
PPT_TITLE = "" # 此 PPT 的标题，变量名必须为 PPT_TITLE

# 你可以在这里定义一些公用的全局变量，例如颜色

# 你可以在这里定义一些可复用的函数，例如每页都会出现的元素（标题、页等）
# 请发挥你的想象力，定义一些合适的函数来复用一些常见的元素
# 绝对禁止只定义一个 x 函数然后所有 make_page_xxx 都使用 x 函数来实现（这样每个页面长得都一模一样，这是绝对禁止的）


# 这里开始制作每页PPT，函数名必须MUST是 make_page_[number] 这种格式，number 从 1 开始，依次递增
# 此代码会交给自动化工作流来处理，如果你不遵守 make_page_[number] 这种格式，将会触发错误
# def make_page_1():
#     slide = new_slide()
#     # 在这里组织第 1 页的所有内容


# def make_page_2():
#     slide = new_slide()
#     # 在这里组织第 2 页的所有内容


# # ... 继续添加 make_page_xxx 函数


# def make_page_n():
#     slide = new_slide()
#     # 在这里组织第 n 页的所有内容


# !!! *** 此区域禁止修改 开始 *** !!!
# 在这里调用所有 make_page_xxx 函数
page_functions = []
for name, obj in list(globals().items()): # 必须使用 list 包裹，否则将发生 RuntimeError: dictionary changed size during iteration
    if name.startswith("make_page_") and callable(obj):
        suffix = name.split("make_page_")[-1] 
        if suffix.isdigit():
            page_functions.append((int(suffix), obj))
page_functions.sort(key=lambda x: x[0])
is_error = False
for page_num, func in page_functions:
    try:
        func()
    except Exception as e:
        print(f"Error in make_page_{page_num}: {e}")
        is_error = True
if is_error:
    raise Exception("制作 PPT 失败")
PRS.save(OUTPUT_PATH) # 保存
# !!! *** 此区域禁止修改 结束 *** !!!
    """
    return content.strip()


def make_prompt(report: str, pptxsdkdoc: str) -> str:
    report = report.strip()
    pptxsdkdoc = pptxsdkdoc.strip()
    today = time.strftime("%Y-%m-%d")
    prompt = f"""
<你的角色>
你是一个专业的经验丰富的python PPT制作专家，擅长使用python代码将调研报告制作成PPT。
</你的角色>

<你的任务>
根据调研报告，编写 main.py 文件，制作一个'报告型'的PPT。目的是为了内部报告，所以内容要详实、逻辑要严谨，能够让人独立看懂。
PPT的受众是行业内的专业人士，他们喜欢信息密度高的PPT，所以你每页PPT的内容要极可能的详细，信息密度高，文字密度高。
如果我给你提供了图片，请挑选合适的图片放到PPT中，图片要恰如其分，不要为了放图片而放图片，如果没有合适的图片，可以不添加任何图片。
请你合理规划 PPT 的页数，页数要与「调研报告」的内容相匹配，「调研报告」的内容多则页数多，内容少则页数少。且最大页数不要超过30页。
PPT的颜色和排版要美观
</你的任务>

<输出格式>
完整的可运行的 main.py 文件，不要任何解释、说明等其他内容
只要输出 main.py 文件！，禁止出现其他内容，禁止用```python包裹代码
</输出格式>

<pptxsdk中可用函数介绍>
{pptxsdkdoc}
</pptxsdk中可用函数介绍>

<main.py文件>
这是整个项目的入口文件，请你完成此文件来实现生成PPT的功能
请遵循main.py文件中的说明
----------以下是main.py文件----------
{get_main_py_init_content()}
----------以上是main.py文件----------
</main.py文件>

<注意事项>
- 你的输出是完整的 main.py 文件内容，禁止出现其他内容，禁止用 ```python 包裹代码
- 要适当添加 fontawesome icon，不要添加太多，也不要太少，要恰到好处
- 如果我给你提供了图片，请挑选合适的图片放到 PPT 中，图片要恰如其分，不要为了放图片而放图片，如果没有合适的图片，可以不添加任何图片。
- 如果你添加了图片，请在图片附近说明图片的来源(但是禁止出现[cite:ID]这样的字样，而是应该从图片内容来推测图片的来源)
- 添加图片时请直接使用图片名称，不要带路径，不要在代码中检测图片是否存在，代码运行在虚拟环境，你根本无需关心工作目录在哪里，只需要使用图片名称即可。
- 一定要注意 python 语法，不要出现语法错误，尤其是用双引号嵌套时，一定要避免！
- PPT 的语言必须与调研报告的语言保持一致，调研报告是中文，则PPT也是中文，调研报告是英文，则PPT也是英文，调研报告是日语，则PPT也是日语
- PPT 的颜色和排版要美观，请从 "推荐配色" 中选择一组合适的配色，并使用该配色来设计 PPT
- 如果你添加了 chart，请在 chart 附近简短解释一下图表的含义，同时在图表标题或者附近说明数据的**单位**

</注意事项>

<你经常犯的错误这次不要再犯了>
- python 中的字符串内容，如果出现双引号中嵌套双引号，请记得加转义字符！！！
- python 中你定义的 颜色全局变量，你使用的时候要把变量名写对！禁止后面的代码使用前面定义的全局变量时写错变量名！
</你经常犯的错误这次不要再犯了>

<推荐配色>
Palette 01: #12355B, #1F6FBA, #7BC6E6, #F4F9FC, #00A6A6, 等其它适合此色卡的配色
Palette 02: #1E1B4B, #5B5FC7, #3A86FF, #A78BFA, #F7F5FF, 等其它适合此色卡的配色
Palette 03: #0B1F3A, #1D4ED8, #4B6B88, #C9A227, #F8FAFC, 等其它适合此色卡的配色
Palette 04: #2563EB, #1F2937, #64748B, #F8FAFC, #06B6D4, 等其它适合此色卡的配色
Palette 05: #0F3B2E, #1E5B4F, #8FAF9D, #D4AF37, #F7F4EC, 等其它适合此色卡的配色
Palette 06: #08111F, #111827, #3B82F6, #22D3EE, #FBBF24, 等其它适合此色卡的配色
Palette 07: #4A102A, #7A1F3D, #B08A99, #E7D8DD, #D6B25E, 等其它适合此色卡的配色
Palette 08: #2F3A4A, #5A6B7D, #9AA7B2, #DCE3E8, #F8FAFB, 等其它适合此色卡的配色
</推荐配色>

<调研报告>
{report}
</调研报告>

<其他信息>
今天是 {today}
</其他信息>

现在请你输出完整的 main.py 文件内容，不要任何解释、说明等其他内容，禁止用```python包裹代码
    """
    return prompt.strip()