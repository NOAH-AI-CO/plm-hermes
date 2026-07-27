import asyncio
import json
import os
import re
import httpx
import base64
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

from agent.bp.pp import parse_and_select
from agent.bp.ppt_image_pipeline import parse_pdf_for_ppt, PptImageAsset

proxy_url = "http://127.0.0.1:7897" 

# 强制设置环境变量
os.environ['http_proxy'] = proxy_url
os.environ['https_proxy'] = proxy_url
os.environ['all_proxy'] = proxy_url

gcp_key_path = "/Users/ivylyx/Code/NoahAgent/noah_agent/gcp_key.json"
if not os.environ.get('GOOGLE_APPLICATION_CREDENTIALS', ''):
    os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = gcp_key_path

os.environ['GOOGLE_CLOUD_PROJECT'] = "noahai-440408"
os.environ['GOOGLE_CLOUD_LOCATION'] = "global"
os.environ['GOOGLE_GENAI_USE_VERTEXAI'] = "true"

from google import genai
from google.genai.types import HttpOptions
import json

_genai_client = None
_httpx_client = None

def get_genai_client():
    global _genai_client
    if _genai_client is None:
        _genai_client = genai.Client(http_options=HttpOptions(api_version="v1"))
    return _genai_client

async def get_httpx_client():
    global _httpx_client
    if _httpx_client is None:
        _httpx_client = httpx.AsyncClient(follow_redirects=True, timeout=httpx.Timeout(30.0))
    return _httpx_client

client = get_genai_client()

# 大纲和内容的 Schema 定义，供 LLM 生成 JSON 时参考使用
outline_schema = {
    "type": "ARRAY",
    "description": "PPT 的完整大纲列表，必须包含：封面、目录、正文、致谢。每项代表一页幻灯片。",
    "items": {
        "type": "OBJECT",
        "properties": {
            "slide_index": {"type": "INTEGER", "description": "页码，从 1 开始连续编号"},
            "title": {"type": "STRING", "description": "该页的标题"},
            "intent": {
                "type": "STRING",
                "description": "排版意图。固定页必须使用：封面、目录、致谢；正文页使用：背景介绍、机制阐述、核心数据展示、结论总结、研究设计、基线特征 等"
            }
        },
        "required": ["slide_index", "title", "intent"]
    }
}

# 单块 Schema：用于 content_schema 的 blocks 数组
block_schema = {
    "type": "OBJECT",
    "properties": {
        "type": {"type": "STRING", "enum": ["text", "image"], "description": "块类型：text=文本块，image=图片块"},
        "label": {"type": "STRING", "description": "可选。块的语义标签，如：左侧要点、右侧数据、主图、机制图。便于排版时区分。"},
        "content": {
            "type": "ARRAY",
            "items": {"type": "STRING"},
            "description": "文本块专用。每条为一段精简文本，支持项目符号。"
        },
        "context": {
            "type": "STRING",
            "description": "图片块专用。该图的画面描述（如：Kaplan-Meier 生存曲线、分子机制图）。"
        },
        "selected_filename": {
            "type": "STRING",
            "description": "图片块专用。从源图片库选中的文件名。若无合适图片可留空。"
        }
    },
    "required": ["type"]
}

content_schema = {
    "type": "OBJECT",
    "properties": {
        "title": {"type": "STRING", "description": "该页的 PPT 标题"},
        "blocks": {
            "type": "ARRAY",
            "items": block_schema,
            "description": "内容块列表。每页 2-6 个块，常规 4 个。块顺序即排版时的逻辑顺序。可混排 text 与 image。"
        }
    },
    "required": ["title", "blocks"]
}

cover_content_schema = {
    "type": "OBJECT",
    "properties": {
        "title": {"type": "STRING", "description": "研究主标题（论文题目或核心结论短语）"},
        "subtitle": {"type": "STRING", "description": "副标题，如作者、机构、会议/期刊、日期等，一行或两行"},
        "paragraphs": {
            "type": "ARRAY",
            "items": {"type": "STRING"},
            "description": "封面次要信息，每项一行。如：作者名单、机构、日期。若 subtitle 已涵盖可留空 []"
        }
    },
    "required": ["title", "subtitle"]
}

system_prompt = """
角色：你是一位顶级医学学术顶会（如 ASCO, ESMO）的视觉主编。你的任务是将【页面内容块】转化为带有绝对坐标的 JSON 排版数据。

【画布物理法则与防穿模机制（极度重要！）】
1. 画布尺寸：宽 (w) 严格固定 10.0 英寸，高 (h) 严格固定 5.625 英寸。绝不允许越界！
2. 标题区域控制（紧凑）：标题应位于页面上方，y 建议 0.20-0.35，h 建议 0.35-0.60。禁止给标题分配过高区域，避免挤压正文和图片。
   - 通用边界：所有 elements 必须满足 x>=0.20, y>=0.15, x+w<=9.80, y+h<=5.45。
   - 通用间距：任何两个 elements（title/text/image）矩形边界的最小间距 >= 0.12 英寸（允许极小浮点误差）。
   - 铺满底部：除 title 外，所有主体元素的最底部 y+h 建议尽量靠近 5.30（允许底部留白 <=0.25）。
   - 标题高度随块数量变化（强制套用）：
     • block_count=2–3：title.h 取 0.40–0.48
     • block_count=4：title.h 取 0.36–0.45
     • block_count=5–6：title.h 取 0.34–0.42
   - 标题下沿到主体的间距：主体元素的 y 必须 >= title.y + title.h + 0.12。
3. 禁止使用水平分割线：不要输出任何“横线”元素（包括用极扁矩形模拟横线）。
4. 卡片高度冗余：若单侧正文超过 60 字，底层的 shape 和 text 的高度必须 >= 2.8。

【多块布局核心规则】每页有 2–6 个内容块（常规 4 块）。必须为输入的每一个内容块输出对应的 elements 元素！块与块严禁重叠，留 0.12 英寸间距。elements 顺序与输入块顺序一致。
【文本视觉规则】text 元素禁止使用外层方形 shape 包裹（不要给 text 增加卡片边框、矩形底框、描边框）。

【多场景动态版式选择（KISS 原则：看意图下菜碟）】
请务必根据传入的【排版意图】和【全局位置】，从以下范式中选择最合适的一种：

范式 A【封面页 / 目录页 / 致谢页 / 宏观总结页】：
- 触发：意图包含“封面”、“背景介绍首页”或“结论总结”。
- 排版：不使用水平分割线。采用大气的中心或偏上居中排版。正文/副标题居中。封面、致谢页可极简；目录页以项目符号或编号列出各章节标题。不需要强行放置图片。

范式 B【核心数据图表优先页】：
- 触发：意图包含“数据展示”、“机制”、“疗效”等，且包含关键图片。
- 排版：顶部紧凑标题 + 主体内容 + 底部脚注。图片 (image) 是绝对主角 (w >= 5.5, h >= 3.5)。正文压缩到侧边栏 (w=3.5 左右)。

范式 C【学术对比卡片页】：
- 触发：意图为“研究设计”、“基线特征”或两组方案对比。
- 排版：顶部紧凑标题。下方使用左右对称的双栏排版。底色使用 #FFFFFF，配上纤细的 #CBD5E1 边框。

范式 D【多块混排】：根据块数量（2–6）灵活布局，如 2 图+2 文可左右或上下分栏，4 块可 2×2 网格。

【视觉基调】
1. 背景色：纯白 #FFFFFF 或超浅临床灰 #F8FAFC。
2. 字体：大标题深蓝 #0F4C81，正文高级冷灰 #334155。

以下为具体输入信息（请仔细阅读全局大纲和当前意图，严格验算坐标）：
{user_content}
"""

location_response_schema = {
  "type": "OBJECT",
  "properties": {
    "slide_index": {
      "type": "INTEGER",
      "description": "幻灯片的页码"
    },
    "elements": {
      "type": "ARRAY",
      "description": "幻灯片上的元素列表。必须为输入的每个内容块输出对应元素（title/text/image）。可含多个 text、多个 image。数组顺序为图层顺序（Z-index）。",
      "items": {
        "type": "OBJECT",
        "properties": {
          "type": {
            "type": "STRING",
            "enum": ["shape", "title", "text", "image"],
            "description": "元素类型"
          },
          "x": {
            "type": "NUMBER",
            "description": "元素左上角的 X 坐标（单位：英寸，范围 0.0 到 10.0）"
          },
          "y": {
            "type": "NUMBER",
            "description": "元素左上角的 Y 坐标（单位：英寸，范围 0.0 到 5.625）"
          },
          "w": {
            "type": "NUMBER",
            "description": "元素宽度（单位：英寸）"
          },
          "h": {
            "type": "NUMBER",
            "description": "元素高度（单位：英寸）"
          },
          "content": {
            "type": "ARRAY",
            "items": {
              "type": "STRING"
            },
            "description": "文本内容。每一项代表一个段落。仅 type 为 title/text 时生效。"
          },
          "font_size": {
            "type": "INTEGER",
            "description": "字体大小（Pt），推荐标题 24-36，正文 14-18"
          },
          "font_color": {
            "type": "STRING",
            "description": "字体颜色的 Hex 值，如 #FFFFFF 或 #333333"
          },
          "is_bold": {
            "type": "BOOLEAN",
            "description": "是否加粗"
          },
          "alignment": {
            "type": "STRING",
            "enum": ["left", "center", "right", "justify"],
            "description": "对齐方式"
          },
          "is_bullet": {
            "type": "BOOLEAN",
            "description": "是否启用项目符号（Bullet points）"
          },
          "line_spacing": {
            "type": "NUMBER",
            "description": "行距倍数，推荐 1.0 或 1.5"
          },
          "shape_type": {
            "type": "STRING",
            "enum": ["rectangle"],
            "description": "形状类型，仅 type 为 shape 时生效。"
          },
          "fill_color": {
            "type": "STRING",
            "description": "填充颜色的 Hex 值，用于 shape"
          },
          "border_color": {
            "type": "STRING",
            "description": "边框颜色的 Hex 值，留空代表无边框"
          },
          "context": {
            "type": "STRING",
            "description": "用于生成或查找图片的提示词描述，仅 type 为 image 时生效。"
          }
        },
        "required": ["type", "x", "y", "w", "h"]
      }
    }
  },
  "required": ["slide_index", "elements"]
}

# 颜色转换辅助函数
def hex_to_rgb(hex_str):
    """将 #FFFFFF 转换为 python-pptx 需要的 RGBColor 对象"""
    if not hex_str or not hex_str.startswith('#'):
        return RGBColor(0, 0, 0) # 默认黑色
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY
}

class PPTGenerator:
    def __init__(self):
        """
        初始化生成器。可以传入自定义模板，默认为 python-pptx 的内置模板。
        """
        self.prs = Presentation()
        # 强制设置为 16:9 比例 (10 inches x 5.625 inches)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        # 可选：用于在 PPT 中插入源 PDF 图片
        self._source_images_dict = {}
        self._source_image_assets: list[PptImageAsset] = []
        self._source_assets_by_filename: dict[str, PptImageAsset] = {}
        self._global_title_font_size: int = 24
        self._title_box_width_in: float = 8.8

    def set_source_images(self, images_dict: dict, assets: list[PptImageAsset] | None = None):
        raw = images_dict or {}
        self._source_images_dict = dict(raw)
        # 为 images/xxx.jpg 形式的 key 添加 xxx.jpg 别名，便于匹配
        for k, v in raw.items():
            if "/" in k and k.endswith(".jpg"):
                base = k.split("/")[-1]
                if base not in self._source_images_dict:
                    self._source_images_dict[base] = v
        self._source_image_assets = assets or []
        self._source_assets_by_filename = {}
        for asset in self._source_image_assets:
            self._source_assets_by_filename[asset.filename] = asset
            if "/" in asset.filename:
                self._source_assets_by_filename[asset.filename.split("/")[-1]] = asset

    def _load_title_font(self, size_pt: int):
        """加载测量标题用字体。若系统字体不可用则返回 None。"""
        try:
            from PIL import ImageFont  # type: ignore
        except Exception:
            return None
        candidates = [
            "/System/Library/Fonts/Helvetica.ttc",
            "/System/Library/Fonts/PingFang.ttc",
            "/Library/Fonts/Arial.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]
        for p in candidates:
            if Path(p).exists():
                try:
                    return ImageFont.truetype(p, size_pt)
                except Exception:
                    continue
        return None

    def _measure_text_width_px(self, text: str, size_pt: int) -> float | None:
        """返回文本在指定字号下的宽度像素。失败时返回 None。"""
        if not text:
            return 0.0
        try:
            from PIL import Image, ImageDraw  # type: ignore
            font = self._load_title_font(size_pt)
            if font is None:
                return None
            img = Image.new("RGB", (4000, 300), "white")
            draw = ImageDraw.Draw(img)
            bbox = draw.textbbox((0, 0), text, font=font)
            return float(max(0, bbox[2] - bbox[0]))
        except Exception:
            return None

    def configure_title_style_from_outline(self, outline: list[dict]) -> None:
        """
        根据提纲中最长标题动态设置全局标题字号，目标是单行不换行。
        使用真实文本宽度测量 + 二分搜索，找可容纳最长标题的最大字号。
        """
        titles = [str(item.get("title", "")).strip() for item in outline if str(item.get("title", "")).strip()]
        if not titles:
            self._global_title_font_size = 24
            return
        longest_title = max(titles, key=len)
        # 以 96 DPI 估算英寸到像素，留 8% 安全边距，减少跨平台渲染误差导致的换行
        available_width_px = self._title_box_width_in * 96.0 * 0.92
        low, high = 12, 36
        best = 18
        # 先尝试真实测量
        measured = self._measure_text_width_px(longest_title, best)
        if measured is not None:
            while low <= high:
                mid = (low + high) // 2
                width_px = self._measure_text_width_px(longest_title, mid)
                if width_px is None:
                    break
                if width_px <= available_width_px:
                    best = mid
                    low = mid + 1
                else:
                    high = mid - 1
        else:
            # Pillow 不可用时的兜底：按字符数估算
            max_len = len(longest_title)
            if max_len <= 12:
                best = 30
            elif max_len <= 18:
                best = 26
            elif max_len <= 24:
                best = 22
            elif max_len <= 30:
                best = 19
            else:
                best = 17
        self._global_title_font_size = max(12, min(36, best))
        print(
            f"   📐 全局标题字号已根据最长标题({len(longest_title)}字)单行适配为: "
            f"{self._global_title_font_size}pt"
        )

    def _parse_source_image_from_context(self, ctx: str) -> str | None:
        """从 context 中解析 SOURCE_IMAGE=<filename>，支持行首或行内出现"""
        s = str(ctx or "").strip()
        if not s:
            return None
        # 在整段 context 中搜索 SOURCE_IMAGE=xxx.jpg（LLM 常将描述与 SOURCE_IMAGE 写在同一行）
        match = re.search(r"SOURCE_IMAGE=([a-zA-Z0-9_\-\.]+\.jpg)", s, re.IGNORECASE)
        if match:
            return match.group(1)
        match = re.search(r"SOURCE_IMAGE=([a-f0-9]{16,}\.jpg)", s)
        if match:
            return match.group(1)
        return None

    def _resolve_image_uri(self, source_filename: str) -> str | None:
        """根据 filename 解析图片 data_uri，支持多种 key 格式"""
        if not source_filename:
            return None
        data_uri = self._source_images_dict.get(source_filename)
        if data_uri:
            return data_uri
        alt_key = f"images/{source_filename}"
        data_uri = self._source_images_dict.get(alt_key)
        if data_uri:
            return data_uri
        for key in self._source_images_dict:
            if key.endswith(source_filename) or source_filename in key:
                return self._source_images_dict[key]
        return None

    def _fit_image_in_box(self, img_w_px: int, img_h_px: int, box_w_in: float, box_h_in: float):
        """Return (w_in, h_in, dx_in, dy_in) to fit image into box preserving aspect ratio."""
        if img_w_px <= 0 or img_h_px <= 0 or box_w_in <= 0 or box_h_in <= 0:
            return box_w_in, box_h_in, 0.0, 0.0

        img_ar = img_w_px / img_h_px
        box_ar = box_w_in / box_h_in

        if img_ar >= box_ar:
            w_in = box_w_in
            h_in = box_w_in / img_ar
            dx_in = 0.0
            dy_in = (box_h_in - h_in) / 2.0
        else:
            h_in = box_h_in
            w_in = box_h_in * img_ar
            dx_in = (box_w_in - w_in) / 2.0
            dy_in = 0.0
        return w_in, h_in, dx_in, dy_in

    # ==========================================
    # 步骤 1：生成提纲 (宏观)
    # ==========================================
    async def generate_outline_from_paper(self, paper_text):
        print("正在从论文中提取 PPT 提纲...")
        prompt = f"""
        你是一位资深的医学联络官(MSL)。请阅读以下医学论文，为其规划一份**完整结构**的 PPT 汇报提纲。

        【必须遵守的页面结构（共 7–9 页）】
        1. 第 1 页 - 封面：intent 必须为「封面」，title 为研究主标题（如论文题目），用于展示研究名称、作者、机构、日期。
        2. 第 2 页 - 目录：intent 必须为「目录」，title 为「目录」或「Agenda」，用于列出后续正文各章节标题。
        3. 第 3–(N-1) 页 - 正文：根据论文内容规划 4–6 页，intent 可为：背景介绍、机制阐述、核心数据展示、研究设计、基线特征、结论总结 等。
        4. 第 N 页 - 致谢：intent 必须为「致谢」，title 为「Thank you」或「致谢」或「Q&A」，作为结尾页。

        【输出要求】
        - slide_index 从 1 开始连续编号。
        - 第一项 intent 必须包含「封面」，第二项必须包含「目录」，最后一项必须包含「致谢」。
        - 正文页的 title 和 intent 需与论文内容匹配，保留核心临床数据（PFS/OS、HR、p值等）的展示。

        论文内容：
        {paper_text}
        """
        response = await client.aio.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config={
                "response_mime_type": "application/json",
                "temperature": 0.2,
                "response_schema": outline_schema
            }
        )
        return json.loads(response.text)

    def _is_special_page(self, intent: str) -> str | None:
        """判断是否为特殊页（封面/目录/致谢），返回意图关键词或 None"""
        intent_lower = (intent or "").strip().lower()
        if "封面" in intent_lower:
            return "封面"
        if "目录" in intent_lower or "agenda" in intent_lower:
            return "目录"
        if "致谢" in intent_lower or "thank" in intent_lower or "q&a" in intent_lower or "qa" in intent_lower:
            return "致谢"
        return None

    async def generate_cover_content_from_paper(self, paper_text: str) -> dict:
        """从论文中提取封面页内容（标题、作者、机构、日期）"""
        print("正在提取封面页内容...")
        prompt = f"""
        请从以下医学论文中提取封面页所需信息。
        输出 JSON：title（研究主标题，如论文题目）、subtitle（副标题，如作者、机构、期刊、年份等，一行概括）、paragraphs（可选的额外信息，如具体作者名单，每项一行，若无则 []）。

        论文内容（摘要或前几段）：
        {paper_text[:8000]}
        """
        response = await client.aio.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config={
                "response_mime_type": "application/json",
                "temperature": 0.2,
                "response_schema": cover_content_schema
            }
        )
        raw = json.loads(response.text)
        paragraphs = [raw["subtitle"]]
        if raw.get("paragraphs"):
            paragraphs.extend(raw["paragraphs"])
        return {
            "title": raw["title"],
            "blocks": [{"type": "text", "label": "副标题", "content": paragraphs}]
        }

    def build_toc_content_from_outline(self, outline: list[dict]) -> dict:
        """根据大纲构建目录页内容（仅列出正文页，排除封面/目录/致谢）"""
        body_titles: list[str] = []
        for item in outline:
            intent = (item.get("intent") or "").strip().lower()
            if "封面" in intent or "目录" in intent or "致谢" in intent or "thank" in intent or "q&a" in intent:
                continue
            title = item.get("title", "").strip()
            if title:
                body_titles.append(title)
        return {
            "title": "目录",
            "blocks": [{"type": "text", "label": "章节列表", "content": body_titles}]
        }

    def build_thanks_content(self) -> dict:
        """构建致谢/Q&A 页的固定内容"""
        return {
            "title": "Thank you",
            "blocks": [{"type": "text", "content": ["Questions?"]}]
        }

    def _normalize_content_to_blocks(self, content_data: dict) -> dict:
        """将内容统一为 blocks 格式。兼容旧格式（封面/目录/致谢的 title+paragraphs）。"""
        if "blocks" in content_data and content_data["blocks"]:
            blocks = content_data["blocks"][:6]  # 最多 6 块
            return {"title": content_data["title"], "blocks": blocks}
        # 旧格式：title + paragraphs + 可选 image
        blocks: list[dict] = []
        paras = content_data.get("paragraphs", [])
        if paras:
            blocks.append({"type": "text", "content": paras})
        img_ctx = content_data.get("image_context", "").strip()
        img_fn = content_data.get("selected_image_filename", "").strip()
        if img_ctx or img_fn:
            blocks.append({"type": "image", "context": img_ctx or "示意图", "selected_filename": img_fn})
        if len(blocks) < 2 and content_data.get("title"):
            blocks.append({"type": "text", "content": [""]})
        return {"title": content_data.get("title", ""), "blocks": blocks}

    # ==========================================
    # 步骤 2：生成单页内容 (微观)
    # ==========================================
    async def generate_content_from_outline(self, outline_item, paper_text):
        print(f"正在生成第 {outline_item['slide_index']} 页的具体内容：{outline_item['title']}...")
        
        # 准备可用的源图片库信息（如果存在）
        image_library_info = ""
        if self._source_image_assets:
            image_library_lines = []
            for asset in self._source_image_assets[:20]:  # 限制数量，避免 prompt 过长
                desc = (asset.description or "").replace("\n", " ").strip()
                dim = ""
                if getattr(asset, "width_px", None) and getattr(asset, "height_px", None):
                    ar = f"{asset.aspect_ratio:.3f}" if getattr(asset, "aspect_ratio", None) else "?"
                    dim = f" ({asset.width_px}x{asset.height_px}px, AR={ar})"
                if desc:
                    image_library_lines.append(f"- {asset.filename}{dim}: {desc[:150]}")
                else:
                    image_library_lines.append(f"- {asset.filename}{dim}: (no description)")
            if image_library_lines:
                image_library_info = f"""
        
        【可用的源 PDF 图片库】:
        {chr(10).join(image_library_lines)}
        
        【重要】每个 type="image" 的块可在 selected_filename 中指定从图片库选择的文件名。可选 0-3 张图，尽量不重复。若无合适图片，selected_filename 留空，但 context 仍需描述预期画面。
        """
        
        prompt = f"""
        你是一位顶级的医学文案专家。我们正在制作一页关于 {outline_item['title']} 的 PPT。
        它的页面意图是：{outline_item['intent']}。
        
        【块数量约束】每页必须输出 2-6 个 blocks，常规为 4 个。根据页面意图灵活组合：
        - 数据展示/机制页：可 2 个 text + 2 个 image，或 1 text + 1 image 等。
        - 基线特征/对比页：可 2 个 text（左右栏）+ 0-2 个 image。
        - 结论页：可 2-3 个 text。
        
        【块要求】text 块的 content 为精简项目符号短句；image 块的 context 描述画面，selected_filename 从图片库选（可选）。必须保留核心临床数据（HR、P值、中位生存期等）。
        
        论文原文：
        {paper_text}
        {image_library_info}
        """
        response = await client.aio.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config={
                "response_mime_type": "application/json",
                "temperature": 0.2,
                "response_schema": content_schema
            }
        )
        return json.loads(response.text)
    
    # ==========================================
    # 步骤 3：JSON 布局生成 (动态接入)
    # ==========================================
    async def generate_slide_json(self, content_data: dict, slide_index: int, intent: str, full_outline_str: str) -> str:
        print(f"   -> 正在计算第 {slide_index} 页 ({intent}) 的排版坐标与视觉 JSON...")
        
        # 统一为 blocks 格式
        norm = self._normalize_content_to_blocks(content_data)
        title = norm["title"]
        blocks = norm["blocks"]
        
        # 构建多块内容描述
        blocks_desc_lines: list[str] = []
        pre_selected_by_block: dict[int, str] = {}  # block_idx -> filename
        for i, blk in enumerate(blocks):
            blk_type = blk.get("type", "text")
            label = blk.get("label", "")
            label_str = f" [{label}]" if label else ""
            if blk_type == "text":
                content = blk.get("content", []) or []
                lines = "\n".join([f"  - {line}" for line in content if line])
                blocks_desc_lines.append(f"  块{i+1} [text]{label_str}:\n{lines or '  (空)'}")
            else:
                ctx = blk.get("context", "").strip() or "示意图"
                fn = blk.get("selected_filename", "").strip()
                blocks_desc_lines.append(
                    f"  块{i+1} [image]{label_str}: context={ctx}"
                    + (f", SOURCE_IMAGE={fn}" if fn else "")
                )
                if fn:
                    pre_selected_by_block[i] = fn
        
        blocks_str = "\n".join(blocks_desc_lines)
        
        # 源图片库
        image_library_lines = []
        for asset in self._source_image_assets[:30]:
            desc = (asset.description or "").replace("\n", " ").strip()
            dim = ""
            if getattr(asset, "width_px", None) and getattr(asset, "height_px", None):
                ar = f"{asset.aspect_ratio:.3f}" if getattr(asset, "aspect_ratio", None) else "?"
                dim = f" ({asset.width_px}x{asset.height_px}px, AR={ar})"
            if desc:
                image_library_lines.append(f"- {asset.filename}{dim}: {desc[:200]}")
            else:
                image_library_lines.append(f"- {asset.filename}{dim}: (no description)")
        image_library_str = "\n".join(image_library_lines) if image_library_lines else "(no source images)"
        
        # 图片绑定规则
        selected_list = list(pre_selected_by_block.values())
        valid_selected = [f for f in selected_list if f in self._source_images_dict]
        if valid_selected:
            bind_rules = "\n".join([f"- 块{k+1} 的 image 必须写 SOURCE_IMAGE={v}" for k, v in pre_selected_by_block.items() if v in self._source_images_dict])
            image_instruction = f"""
        【已选中的图片（按块绑定）】:
        {bind_rules}
        """
            for f in valid_selected:
                print(f"   📌 使用步骤1已选中的图片: {f}")
        else:
            image_instruction = f"""
        【可用的源 PDF 图片库】:
        {image_library_str}
        【规则】每个 type="image" 的元素需在 context 中写 SOURCE_IMAGE=<filename>，filename 来自上图库。不同 image 尽量选不同文件。
        """
        
        user_input = f"""
        【全局 PPT 结构】: 
        {full_outline_str}
        
        【当前处理页面】: 
        幻灯片页码：{slide_index}
        排版意图：{intent}
        
        【页面内容 - 多块布局（共 {len(blocks)} 块）】:
        标题：{title}
        块数量 block_count={len(blocks)}
        内容块（按块排版，输出对应的 title/text/image 元素，块间不重叠）：
        {blocks_str}
        {image_instruction}
        文献参考：Modi S, et al. N Engl J Med. 2022.
        """

        response = await client.aio.models.generate_content(
            model="gemini-3-flash-preview",
            contents=system_prompt.format(user_content=user_input),
            config={
                "response_mime_type": "application/json",
                "temperature": 0.2,
                "response_schema": location_response_schema
            },
        )
        return response.text

    # ==========================================
    # 步骤 4：添加单页幻灯片 (不再直接保存)
    # ==========================================
    def add_slide_from_json(self, json_data):
        """解析 JSON 数据并添加到当前 Presentation 实例中"""
        slide_data = json.loads(json_data)
        is_first_page = int(slide_data.get("slide_index", 0) or 0) == 1
        
        # 始终添加纯空白版式的幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        text_boxes: list[tuple[float, float, float, float]] = []
        for el in slide_data.get("elements", []):
            if el.get("type") in ("title", "text"):
                x = float(el.get("x", 0))
                y = float(el.get("y", 0))
                w = float(el.get("w", 0))
                h = float(el.get("h", 0))
                if w > 0 and h > 0:
                    text_boxes.append((x, y, x + w, y + h))

        for element in slide_data.get("elements", []):
            el_type = element.get("type")
            
            left = Inches(element.get("x", 0))
            top = Inches(element.get("y", 0))
            width = Inches(element.get("w", 1))
            height = Inches(element.get("h", 1))

            # --- 处理文本和标题 ---
            if el_type in ["title", "text"]:
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.word_wrap = True if (el_type == "text" or is_first_page) else False
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0

                raw_content = element.get("content", "")
                paragraphs_text = raw_content if isinstance(raw_content, list) else raw_content.split('\n')
                
                # 防御层：字数动态降维算法
                total_chars = sum(len(str(line)) for line in paragraphs_text)
                base_font_size = element.get("font_size", 14)
                if el_type == "title":
                    # 第 1 页不限制标题字号；其余页使用全局字号确保单行风格一致
                    if not is_first_page:
                        base_font_size = self._global_title_font_size
                
                if el_type == "text":
                    if total_chars > 60 and base_font_size > 12:
                        base_font_size = 12
                    if total_chars > 100 and base_font_size > 10:
                        base_font_size = 11

                for idx, text_line in enumerate(paragraphs_text):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    if element.get("is_bullet", False):
                        text_line = text_line.strip()
                        if not text_line.startswith("•") and not text_line.startswith("-"):
                            text_line = f"• {text_line}"
                    
                    p.text = text_line
                    p.space_after = Pt(0)
                    p.space_before = Pt(0)
                    
                    if "line_spacing" in element:
                        p.line_spacing = element["line_spacing"]
                    
                    p.font.size = Pt(base_font_size)
                    p.font.bold = element.get("is_bold", False)
                    
                    if "font_color" in element:
                        p.font.color.rgb = hex_to_rgb(element["font_color"])
                        
                    if "alignment" in element:
                        p.alignment = ALIGN_MAP.get(element["alignment"], PP_ALIGN.LEFT)

            # --- 处理纯色形状 ---
            elif el_type == "shape":
                # 仅保留全页背景形状；禁止用于包裹 text 的方形卡片
                x1 = float(element.get("x", 0))
                y1 = float(element.get("y", 0))
                x2 = x1 + float(element.get("w", 0))
                y2 = y1 + float(element.get("h", 0))
                is_full_page_bg = (
                    x1 <= 0.05 and y1 <= 0.05 and
                    float(element.get("w", 0)) >= 9.8 and
                    float(element.get("h", 0)) >= 5.4
                )
                if not is_full_page_bg:
                    overlaps_text = False
                    for tx1, ty1, tx2, ty2 in text_boxes:
                        ix = min(x2, tx2) - max(x1, tx1)
                        iy = min(y2, ty2) - max(y1, ty1)
                        if ix > 0.01 and iy > 0.01:
                            overlaps_text = True
                            break
                    if overlaps_text:
                        continue
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                if element.get("fill_color"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(element["fill_color"])
                else:
                    shape.fill.background()
                    
                if element.get("border_color"):
                    shape.line.color.rgb = hex_to_rgb(element["border_color"])
                else:
                    shape.line.fill.background() 

            # --- 处理图片占位符 ---
            elif el_type == "image":
                ctx = element.get('context', '') or ''
                source_filename = self._parse_source_image_from_context(ctx)

                # 尝试匹配图片（支持多种 key 格式）
                data_uri = None
                if source_filename:
                    data_uri = self._resolve_image_uri(source_filename)
                    if not data_uri:
                        print(f"   ❌ 未找到匹配图片: {source_filename}")
                        print(f"   📦 可用图片数量: {len(self._source_images_dict)}")
                        if self._source_images_dict:
                            sample_keys = list(self._source_images_dict.keys())[:3]
                            print(f"   📋 示例 keys: {sample_keys}")

                if data_uri:
                    try:
                        # data_uri expected: "data:image/jpeg;base64,...."
                        if "," in data_uri:
                            _header, b64 = data_uri.split(",", 1)
                        else:
                            b64 = data_uri
                        img_bytes = base64.b64decode(b64)
                        # Preserve aspect ratio: fit image into the (x,y,w,h) box and center it.
                        img_w_px = img_h_px = None
                        try:
                            from PIL import Image  # type: ignore
                            with Image.open(BytesIO(img_bytes)) as im:
                                img_w_px, img_h_px = im.size
                        except Exception:
                            img_w_px = img_h_px = None

                        if img_w_px and img_h_px:
                            box_w_in = float(element.get("w", 1))
                            box_h_in = float(element.get("h", 1))
                            w_in, h_in, dx_in, dy_in = self._fit_image_in_box(img_w_px, img_h_px, box_w_in, box_h_in)
                            slide.shapes.add_picture(
                                BytesIO(img_bytes),
                                Inches(float(element.get("x", 0)) + dx_in),
                                Inches(float(element.get("y", 0)) + dy_in),
                                width=Inches(w_in),
                                height=Inches(h_in),
                            )
                        else:
                            # Fallback: provide only width (pptx will keep aspect ratio)
                            slide.shapes.add_picture(BytesIO(img_bytes), left, top, width=width)
                        print(f"   ✅ 成功插入图片: {source_filename}")
                    except Exception as e:
                        # fallback to placeholder if decode/insert fails
                        print(f"   ❌ 图片插入失败: {e}")
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                        shape.text = f"[ 图片区 - 插入失败 ]\n{ctx}"
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = hex_to_rgb("#F1F5F9")
                        shape.line.color.rgb = hex_to_rgb("#CBD5E1")
                        for p in shape.text_frame.paragraphs:
                            p.alignment = PP_ALIGN.CENTER
                            p.font.color.rgb = hex_to_rgb("#64748B")
                            p.font.size = Pt(12)
                else:
                    # 未找到匹配的图片，显示占位符
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                    shape.text = f"[ 图片区 ]\n{ctx}"
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb("#F1F5F9")
                    shape.line.color.rgb = hex_to_rgb("#CBD5E1")
                    for p in shape.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = hex_to_rgb("#64748B")
                        p.font.size = Pt(12)

        print(f"   ✅ 第 {slide_data.get('slide_index', 'X')} 页幻灯片渲染并装载完成！")

    # ==========================================
    # 步骤 5：最终保存 PPT 文件
    # ==========================================
    def save_ppt(self, output_filepath):
        self.prs.save(output_filepath)
        print(f"\n🎉 全部页面处理完毕！完整演示文稿已成功保存至: {output_filepath}")


# ==========================================
# 测试链路 (Demo Pipeline) - 全自动化多页生成
# ==========================================
async def main():
    import os
    
    generator = PPTGenerator()

    # 1. 读取并解析论文
    pdf_path = '/Users/ivylyx/Code/NoahAgent/noah_agent/agent/bp/s41467-023-42811-4.pdf'
    paper_text_list, images_dict, image_assets = await parse_pdf_for_ppt(pdf_path, detailed=2)
    generator.set_source_images(images_dict, image_assets)
    # 2. 宏观拆解：论文 -> 提纲
    print("[Global Step] 正在让大模型通读全篇长文并规划整体大纲...")
    massive_paper_text = "\n".join(paper_text_list)
    outline = await generator.generate_outline_from_paper(massive_paper_text)
    print(f"✅ 提纲生成完毕，共规划了 {len(outline)} 页：")
    for item in outline:
        print(f"  - 第{item['slide_index']}页: {item['title']} ({item['intent']})")
    generator.configure_title_style_from_outline(outline)

    print("\n" + "="*50)
    print("🚀 开始逐页生成内容与排版流水作业")
    print("="*50)

    # 提取一个极简的全局大纲字符串，透传给排版引擎（KISS 原则）
    full_outline_str = " -> ".join([f"第{item['slide_index']}页({item['intent']})" for item in outline])

    # 3. 循环引擎：逐页提取内容 -> 算坐标 -> 存入内存
    for item in outline:
        current_index = item['slide_index']
        current_intent = item['intent']
        print(f"\n▶ 正在制作第 {current_index} 页: {item['title']}")
        
        # 3.1 根据页面类型提取内容
        special_type = generator._is_special_page(current_intent)
        if special_type == "封面":
            content_data = await generator.generate_cover_content_from_paper(massive_paper_text)
        elif special_type == "目录":
            content_data = generator.build_toc_content_from_outline(outline)
        elif special_type == "致谢":
            content_data = generator.build_thanks_content()
        else:
            content_data = await generator.generate_content_from_outline(item, massive_paper_text)
        
        # 3.2 动态计算排版 JSON (传入 current_intent 和 full_outline_str)
        slide_layout_json = await generator.generate_slide_json(
            content_data=content_data, 
            slide_index=current_index, 
            intent=current_intent,
            full_outline_str=full_outline_str
        )
        
        # 3.3 物理渲染并添加到内存中的幻灯片列表
        generator.add_slide_from_json(slide_layout_json)

    # 4. 最终落盘：将内存中的所有幻灯片保存为一个完整的文件
    output_filename = "/Users/ivylyx/Code/NoahAgent/noah_agent/agent/bp/example.pptx"
    generator.save_ppt(output_filename)

if __name__ == "__main__":
    asyncio.run(main())